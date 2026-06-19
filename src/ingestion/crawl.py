"""Crawl a corpus directory and emit manifest .csv"""
from pathlib import Path
import pandas as pd
from datetime import datetime
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from pptx.enum.shapes import MSO_SHAPE_TYPE

CORPUS_ROOT = Path('/Users/yanghanghuang/Desktop/PhD Related/2024 Summer Poster/')
OUT = Path("manifest.csv")

# --- shared helper -------------------------------------------------------
# _count_images: count images in an *already-opened* python-docx Document.
# Single source of truth — both stat_docx (here) and load_docx_1b (loaders.py)
# call this so we don't re-open the same file twice.
def _count_images(doc) -> int:
    return sum(
        1 for rel in doc.part.rels.values()
        if "image" in rel.target_ref          # (easy) same idiom as before
    )


def _count_pptx_images(prs) -> int:
    """Count picture shapes across all slides of an already-opened Presentation."""
    return sum(
        1 for slide in prs.slides for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )


def _count_pdf_images(reader) -> int:
    """Count embedded /Image XObjects across all pages of an already-opened PdfReader."""
    count = 0
    for page in reader.pages:
        resources = page.get("/Resources")
        if resources and "/XObject" in resources:
            xobj = resources["/XObject"].get_object()
            count += sum(
                1 for obj in xobj.values()
                if obj.get_object().get("/Subtype") == "/Image"
            )
    return count


# --- per-filetype statters ----------------------------------------------
# Walk through the corpus directory and find all .docx and .pptx files
def stat_docx(p: Path) -> dict:
    doc = Document(p)
    text = "\n".join([para.text for para in doc.paragraphs])
    return {
        "word_count": len(text.split()),
        "slide_count": None,
        "image_count": _count_images(doc),            # (think) which helper to call?
        "has_speaker_notes": None,
    }

def stat_pptx(p: Path) -> dict:
    prs = Presentation(p)
    words, notes_flag = 0, False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                words += sum(len(r.text.split())
                             for para in shape.text_frame.paragraphs
                             for r in para.runs)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes_flag = True
            words += len(slide.notes_slide.notes_text_frame.text.split())
    return {
        "word_count": words,
        "slide_count": len(prs.slides),
        "image_count": _count_pptx_images(prs),
        "has_speaker_notes": notes_flag,
    }

def stat_pdf(p: Path) -> dict:
    reader = PdfReader(str(p))
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    return {
        "word_count": len(text.split()),
        "slide_count": None,
        "page_count": len(reader.pages),
        "image_count": _count_pdf_images(reader),
        "has_speaker_notes": None,
    }

STATTERS = {".docx": stat_docx, ".pptx": stat_pptx, ".pdf": stat_pdf}

def crawl(root: Path) -> pd.DataFrame:
    rows = []
    for p in root.rglob("*"):
        if not p.is_file() or p.suffix.lower() not in STATTERS:
            continue
        base = {
            "path": str(p.resolve()),
            "type": p.suffix.lower().lstrip("."),
            "size_kb": round(p.stat().st_size / 1024, 1),
            "modified": datetime.fromtimestamp(p.stat().st_mtime).date().isoformat(),
        }
        try:
            base.update(STATTERS[p.suffix.lower()](p))
            base["error"] = None
        except Exception as e:
            base.update(word_count=None, slide_count=None, page_count=None,
                        image_count=None, has_speaker_notes=None)
            base["error"] = repr(e)
        rows.append(base)
    return pd.DataFrame(rows)

if __name__ == "__main__":
    df = crawl(CORPUS_ROOT)
    df.to_csv(OUT, index=False)
    print(f"{len(df)} files -> {OUT}")
    print(df.groupby("type").agg(
        n=("path", "count"),
        total_words=("word_count", "sum"),
        total_slides=("slide_count", "sum"),
        total_pages=("page_count", "sum"),
        total_images=("image_count", "sum"),
    ))