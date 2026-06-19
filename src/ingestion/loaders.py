"""Ingestion loaders for approaches unstructured (1A), native (1B), and vision (1C)."""

from pathlib import Path
import io
import json
import re
from typing import Iterator

from docx import Document                           # pip install python-docx
from pptx import Presentation                       # pip install python-pptx
from pypdf import PdfReader                         # pip install pypdf

from .schema import ParsedDocument, ParsedSection
from .crawl import _count_images, _count_pptx_images, _count_pdf_images  # shared helpers


def _normalize_approach(approach: str) -> str:
	"""Normalize user-facing approach aliases to internal keys."""
	key = approach.strip().lower()
	aliases = {
		"1a": "1a", "unstructured": "1a", "approach_1a": "1a",
		"1b": "1b", "native": "1b",       "approach_1b": "1b",
		"1c": "1c", "vision": "1c",       "multimodal": "1c", "approach_1c": "1c",
	}
	if key not in aliases:
		raise ValueError(f"Unknown approach: {approach}")
	return aliases[key]


# ---------------------------------------------------------------------------
# Approach unstructured (1A)
# ---------------------------------------------------------------------------
# Three loaders, each split into two scaffolds (level 3 — fewer micro-hints).
# Common pattern across all three:
#   1. _iter_*_unstructured_sections  — walk the flat element list, map categories, extract metadata
#   2. load_*_unstructured            — call partition_*, count Image elements, orchestrate
#
# Key API facts (read before starting):
#   • partition_docx/pptx/pdf(filename=str(path)) → list[Element]
#   • el.category   — e.g. "Title", "NarrativeText", "ListItem", "Image", "Table"
#   • el.text       — extracted text string (may be empty for images)
#   • el.metadata.page_number — 1-based page/slide number (None if unavailable)


def _unstructured_section_label(category: str) -> str:
	"""Map an unstructured element category to a single-word section label.

	"Title" → "heading"   (matches the 1B convention)
	anything else → "normal"
	"""
	return "heading" if category == "Title"  else "normal"           # (easy) — which category string?


# ---- load_docx_unstructured ------------------------------------------------

# ---- Scaffold 1 of 2 -------------------------------------------------------
def _iter_docx_unstructured_sections(elements) -> Iterator[ParsedSection]:
	"""Walk a flat unstructured element list from a .docx; yield one ParsedSection per text element.

	Skip Image elements (counted separately in the loader).
	Track the most recent Title so body paragraphs know their heading context.
	"""
	last_title: str = ""

	for el in elements:
		if el.category == "Image":                                    # (easy) which category to skip?
			continue
		text = el.text.strip()
		if not text:
			continue

		label = _unstructured_section_label(el.category)

		if label == "heading":
			last_title = text                                      # (easy) update the context tracker
			yield ParsedSection(text=text, section=label,
			                    heading_path=[text])                 # (think) same convention as _iter_doc_sections
		else:
			yield ParsedSection(text=text, section=label,
			                    heading_path=[last_title] if last_title else [])                 # (think) inherit from last_title


# ---- Scaffold 2 of 2 -------------------------------------------------------
def load_docx_unstructured(path: Path) -> ParsedDocument:
	"""Approach unstructured (1A): richer element typing, no manual heading stack."""
	from unstructured.partition.docx import partition_docx   # pip install unstructured

	elements = partition_docx(filename=str(path))                       # (easy) wants a str, not a Path
	image_count = sum(1 for el in elements if el.category == "Image") # (think) which category string?
	sections = list(_iter_docx_unstructured_sections(elements))

	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="docx",
		title=path.stem,
		sections=sections,
		meta={"image_count": image_count},
	)


# ---- load_pptx_unstructured ------------------------------------------------

# ---- Scaffold 1 of 2 -------------------------------------------------------
def _iter_pptx_unstructured_sections(elements) -> Iterator[ParsedSection]:
	"""Walk unstructured PPTX elements; yield one section per non-Image element.

	Unlike docx, PPTX has no nested heading hierarchy — use the last Title as
	heading context and carry the slide number from el.metadata.page_number.
	"""
	last_title: str = ""

	for el in elements:
		if el.category == "Image":
			continue
		text = el.text.strip()
		if not text:
			continue

		# unstructured stores the slide number here (may be None)
		slide_idx = el.metadata.page_number if el.metadata else None                                           # (think) what attribute path?

		label = _unstructured_section_label(el.category)

		if label == "heading":
			last_title = text
			yield ParsedSection(text=text, section=label,
			                    slide_idx=slide_idx,
			                    heading_path=[text])
		else:
			yield ParsedSection(text=text, section=label,
			                    slide_idx=slide_idx,
			                    heading_path=[last_title] if last_title else [])


# ---- Scaffold 2 of 2 -------------------------------------------------------
def load_pptx_unstructured(path: Path) -> ParsedDocument:
	"""Approach unstructured (1A): per-element typing + slide metadata."""
	from unstructured.partition.pptx import partition_pptx

	elements = partition_pptx(filename = str(path))                                                # (easy) mirror load_docx_unstructured
	image_count = sum(1 for el in elements if el.category == "Image") # (easy) same pattern
	sections = list(_iter_pptx_unstructured_sections(elements))

	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="pptx",
		title=path.stem,
		sections=sections,
		meta={"image_count": image_count},
	)


# ---- load_pdf_unstructured -------------------------------------------------

# ---- Scaffold 1 of 2 -------------------------------------------------------
def _iter_pdf_unstructured_sections(elements) -> Iterator[ParsedSection]:
	"""Walk unstructured PDF elements; yield one section per non-Image element.

	PDFs have no reliable heading hierarchy; treat "Title" elements as headings anyway.
	Carry the page number from el.metadata.page_number.
	"""
	for el in elements:
		if el.category == "Image":
			continue
		text = el.text.strip()
		if not text:
			continue

		page_idx = el.metadata.page_number if el.metadata else None                                            # (easy) mirror _iter_pptx_unstructured_sections

		label = _unstructured_section_label(el.category)

		yield ParsedSection(
			text=text,
			section=label,
			page_idx=page_idx,
		)


# ---- Scaffold 2 of 2 -------------------------------------------------------
def load_pdf_unstructured(path: Path) -> ParsedDocument:
	"""Approach unstructured (1A): page-level element extraction."""
	from unstructured.partition.pdf import partition_pdf

	elements = partition_pdf(filename = str(path), strategy='hi_res')                                                # (easy)
	image_count = sum(1 for el in elements if el.category == "Image") # (easy)
	sections = list(_iter_pdf_unstructured_sections(elements))

	# Count how many distinct page numbers appear across ALL elements (including Images).
	# Hint: set comprehension over metadata.page_number; filter out None values.
	page_count = len({el.metadata.page_number for el in elements if el.metadata and el.metadata.page_number is not None})                                              # (stretch)

	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="pdf",
		title=path.stem,
		sections=sections,
		meta={"image_count": image_count, "page_count": page_count},
	)


# ---------------------------------------------------------------------------
# Approach native (1B)
# ---------------------------------------------------------------------------

# load_docx_native: extract structured text sections from a .docx file.
# Input:  path — pathlib.Path pointing to a .docx file
# Output: ParsedDocument with heading/paragraph sections and image_count in meta
#
# Split into 3 small scaffolds — complete them in order:
#   1. _heading_level      — turn "heading 2" into int 2
#   2. _iter_doc_sections  — yield ParsedSections, tracking heading ancestry
#   3. load_docx_native     — orchestration; reuses stat_docx from crawl.py


# ---- Scaffold 1 of 3 ------------------------------------------------------
def _heading_level(style: str) -> int:
	"""Turn a docx style name like 'heading 2' into an int depth (default 1)."""

	# Step 1: split the (already-lowercased) style on whitespace into tokens
	#         e.g. "heading 2" -> ["heading", "2"]
	parts = style.split()                                       # (easy)

	# Step 2: grab the last token, but only if it's all digits.
	#         str has a method that returns True only when every char is 0-9.
	last = parts[-1] if parts and parts[-1].isdigit() else "1"    # (think)

	# Step 3: convert and return
	return int(last)                                          # (easy)


# ---- Scaffold 2 of 3 ------------------------------------------------------
def _iter_doc_sections(doc) -> Iterator[ParsedSection]:
	"""Walk a docx; yield one ParsedSection per non-blank paragraph.

	Headings update a stack so every yielded section knows its ancestry path.
	"""

	heading_stack: list[str] = []

	for para in doc.paragraphs:
		text = para.text.strip()
		if not text:
			continue

		style = (para.style.name or "").lower()
		# What word is always present in heading style names ("heading 1", "heading 2", ...)?
		is_heading = "heading" in style                           # (easy)

		if is_heading:
			level = _heading_level(style)
			# Trim the stack so it only contains ancestors of THIS heading.
			# Example: level=2, stack=['Intro','Method','Sub']  ->  ['Intro']
			# Slice math: how many of the existing items are ancestors of a level-`level` heading?
			heading_stack = heading_stack[:level-1]               # (stretch)
			heading_stack.append(text)
			yield ParsedSection(text=text, section="heading",
			                    heading_path=list(heading_stack))
		else:
			# What single-word section label fits non-heading paragraphs?
			yield ParsedSection(text=text, section="normal",     # (easy)
			                    heading_path=list(heading_stack))


# ---- Scaffold 3 of 3 ------------------------------------------------------
def load_docx_native(path: Path) -> ParsedDocument:
	"""Approach native (1B): python-docx parsing with heading-stack metadata.

	Image counting is delegated to crawl._count_images (single source of truth).
	The helper takes an already-opened Document so we don't re-parse the file.
	"""

	# Step 1: open and collect sections (heading-stack lives in the helper)
	doc = Document(path)
	sections = list(_iter_doc_sections(doc))

	# Step 2: image count on our already-open doc (no re-parse)
	image_count = _count_images(doc)

	# Step 3: assemble the ParsedDocument and return
	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="docx",
		title=path.stem,
		sections=sections,
		meta={"image_count": image_count},
	)


# load_pptx_native: extract structured text sections from a .pptx file.
# Input:  path — pathlib.Path pointing to a .pptx file
# Output: ParsedDocument with one section per non-empty slide.
#
# Two scaffolds (level 3 — fewer hints, more design weight per blank):
#   1. _iter_slide_sections — yield ParsedSection per slide (title, body, notes)
#   2. load_pptx_native     — orchestration; counts pictures across the deck

# ---- Scaffold 1 of 2 ------------------------------------------------------
def _iter_slide_sections(prs) -> Iterator[ParsedSection]:
	"""Walk a Presentation; yield one ParsedSection per non-empty slide.

	Each yielded section sets:
	  - text          = body text (concatenated from non-title text frames)
	  - section       = a single-word label
	  - slide_idx     = 1-based slide number
	  - heading_path  = [slide title] when the slide has one (or [] otherwise)
	  - speaker_notes = the slide's notes pane text, or None
	"""

	# Why start=___ ? slides in PowerPoint are 1-indexed in the UI; align with that.
	for idx, slide in enumerate(prs.slides, start=1):                        # (think)

		# Title (may be None for layouts without a title placeholder)
		title_shape = slide.shapes.title
		title = title_shape.text_frame.text.strip() if title_shape else ""

		# Body text — every text frame on the slide except the title's
		body_parts: list[str] = []
		for shape in slide.shapes:
			if shape == title_shape:
				continue
			# Defensive: not every shape has a text frame (pictures, charts, etc.)
			if shape.has_text_frame:                                              # (easy)
				t = shape.text_frame.text.strip()
				if t:
					body_parts.append(t)
		body = "\n".join(body_parts)

		# Speaker notes — only present if the slide actually has a notes_slide.
		# Use has_notes_slide (boolean, no side effect) NOT notes_slide
		# (which lazily creates a notes_slide if absent — mutates the deck).
		notes = None
		if slide.has_notes_slide:
			notes = slide.notes_slide.notes_text_frame.text.strip() or None

		# Skip slides that are completely empty (no title, no body, no notes)
		if not (title or body or notes):
			continue

		yield ParsedSection(
			text=body or title,
			section="slide",
			slide_idx=idx,
			heading_path=[title] if title else [],
			speaker_notes=notes,
		)


# ---- Scaffold 2 of 2 ------------------------------------------------------
def load_pptx_native(path: Path) -> ParsedDocument:
	"""Approach native (1B): python-pptx parsing, one section per slide."""

	prs = Presentation(path)
	sections = list(_iter_slide_sections(prs))

	image_count = _count_pptx_images(prs)

	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="pptx",
		title=path.stem,
		sections=sections,
		meta={"image_count": image_count},
	)

# load_pdf_native: extract one section per non-empty page of a .pdf file.
# Input:  path — pathlib.Path pointing to a .pdf file
# Output: ParsedDocument with one section per page that has extractable text
#
# Two scaffolds (level 3 — inline hints, difficulty tags):
#   1. _iter_pdf_sections — yield ParsedSection per non-empty page
#   2. load_pdf_native    — orchestration; counts embedded /Image XObjects

# ---- Scaffold 1 of 2 ------------------------------------------------------
def _iter_pdf_sections(reader) -> Iterator[ParsedSection]:
	"""Walk a PdfReader; yield one ParsedSection per page that has extractable text.

	PDFs have no native heading/section structure (without TOC parsing), so we
	use page boundaries instead. Image-only / scanned pages will produce no text
	and should be skipped here — flag them upstream as needing OCR.
	"""

	# Why start=___ ? PDF viewers number pages starting at 1, not 0 — match the UI.
	for idx, page in enumerate(reader.pages, start=1):                          # (easy)

		# extract_text() can return None on image-only pages.
		# Pick a short-circuit operator that yields the SECOND value when the first is None,
		# so .strip() always runs on a string.
		raw = page.extract_text()
		text = (raw or "").strip()                                               # (think)
		if not text:
			continue

		yield ParsedSection(
			text=text,
			# Single-word section label, matching docx ("paragraph") / pptx ("slide") convention.
			section="page",                                                       # was "pdf_paragraph" — name the unit you're yielding
			# Which loop variable carries the current page number?
			page_idx=idx,                                                         # (easy)
		)


# ---- Scaffold 2 of 2 ------------------------------------------------------
def load_pdf_native(path: Path) -> ParsedDocument:
	"""Approach native (1B): pypdf parsing, one section per non-empty page."""

	# pypdf's PdfReader is picky about path types — pass a string, not a Path object.
	# There's a one-word built-in that converts almost anything to a string.
	reader = PdfReader(str(path))                                                   # (easy)

	# Call the iterator helper you just wrote above — exhaust it into a list.
	sections = list(_iter_pdf_sections(reader))                                          # (easy)

	image_count = _count_pdf_images(reader)

	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="pdf",
		title=path.stem,
		sections=sections,
		meta={
			"image_count": image_count,
			# pypdf exposes pages as an indexable sequence. Use the obvious built-in to count it.
			"page_count":  len(reader.pages),                                           # (easy)
		},
	)


# ---------------------------------------------------------------------------
# Approach vision (1C)
# ---------------------------------------------------------------------------
# Two patterns:
#   Pattern A (docx, pptx): native text extraction + vision descriptions of embedded images.
#   Pattern B (pdf):        native text per page; full-page vision fallback on thin pages.
#
# vision_client interface (wired up in your application layer):
#   vision_client(image_bytes: bytes, prompt: str) -> str
#   Pass None to disable vision — all three loaders degrade gracefully to native (1B) behaviour.
#
# Extra dependency for Pattern B only:
#   pip install pdf2image pillow   (+ install the poppler system package)
#
# ─── Scaffold rules for this section ────────────────────────────────────────
# Per the new rule in AI_Coding_Teacher.md:
#   • Algorithmic blanks (slice math, control flow, composition) are KEPT.
#   • Trivia (magic strings, library API attributes, repeated patterns) is GIVEN.
#   • Repeated patterns across docx/pptx/pdf are blanked ONCE — the rest are
#     given as worked examples so you study the pattern instead of typing it
#     three times.
#
# Total blanks in this section: 4 algorithmic + 3 trivia (each tagged inline).
# ─────────────────────────────────────────────────────────────────────────────

_VISION_PROMPT = (
	"Describe this image concisely. "
	"If it shows a chart, table, diagram, or experimental result, "
	"summarize the key data or finding it presents."
)
_VISION_THRESHOLD = 20   # words per page below which Pattern B triggers vision

# Caption pattern: matches "Figure 1", "Fig. 3", "fig 12", "Table 4", etc.
#   \b           = word boundary (so "configure" doesn't match)
#   (?:figure|fig\.?|table)  = non-capturing group; .? = optional period after "fig"
#   \s+\d        = whitespace then a digit (to require an actual number)
# False positives we accept: "List of Figures" pages, body refs to figures elsewhere.
# The conjunction with _page_has_image_xobject below filters most of these out.
_FIGURE_CAPTION_RE = re.compile(r"\b(?:figure|fig\.?|table)\s+\d", re.IGNORECASE)


# ---- helper ----------------------------------------------------------------
# BytesIO acts as an in-memory file: PIL writes into it, we pull bytes out via
# .getvalue() (given here because that's the same idiom for any BytesIO use).
def _pil_to_bytes(img) -> bytes:
	"""Encode a PIL Image as bytes suitable for a vision API call."""
	buf = io.BytesIO()
	# Hint: PIL's save() takes a format string naming the image encoder to use.
	# Background: vision APIs accept JPEG and PNG; JPEG is the universal default
	#             (smaller payload, lossy — fine for figures/photos, not for
	#              line art or diagrams where PNG preserves crisper edges).
	# Answer: "JPEG"
	img.save(buf, format="________")               # (easy) trivia — see comment above
	return buf.getvalue()


# ---- load_docx_vision  (Pattern A) -----------------------------------------
# This is where the new section-naming concept lives.
# After this, pptx and pdf reuse the pattern with their own givens.

# ---- Scaffold 1 of 2 -------------------------------------------------------
def _iter_docx_vision_sections(doc, vision_client) -> Iterator[ParsedSection]:
	"""Native text sections first, then one section per embedded image (vision-described).

	If vision_client is None, image sections are skipped (graceful degradation to native).
	"""
	yield from _iter_doc_sections(doc)          # all native text sections first

	if vision_client is None:
		return

	for rel in doc.part.rels.values():
		if "image" not in rel.target_ref:
			continue
		# Hint: docx exposes raw image bytes via the rel's target_part attribute path.
		# Background: .blob is the python-docx convention for raw byte content on any
		#             "part" object — the same naming shows up in python-pptx for
		#             picture shapes (you'll see shape.image.blob below).
		# Answer: rel.target_part.blob
		image_bytes = rel.target_part.blob      # (easy) trivia — see comment above
		description = vision_client(image_bytes, _VISION_PROMPT)
		if description.strip():
			yield ParsedSection(
				text=description,
				# What single-word label fits a vision-described unit?
				# Convention so far: name the unit you're yielding.
				#   docx text → "paragraph" / "heading"
				#   pptx slide → "slide"
				#   pdf page → "page"
				#   docx/pptx image → ???
				section="doc_ppt_image",                   # (easy) — see HINTS at bottom
				meta={"source": "vision"},
			)


# ---- Scaffold 2 of 2 -------------------------------------------------------
def load_docx_vision(path: Path, vision_client=None) -> ParsedDocument:
	"""Approach vision (1C): native docx text + vision descriptions of embedded images."""
	doc = Document(path)
	sections = list(_iter_docx_vision_sections(doc, vision_client))
	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="docx",
		title=path.stem,
		sections=sections,
		meta={
			"image_count":    _count_images(doc),
			# How many sections used vision? Compose a generator with a filter
			# clause that inspects each section's meta dict.
			# (This is the gold-standard "comprehension + filter" idiom — once
			# you write it here, the pptx and pdf loaders below get it for free.)
			"vision_applied": sum(1 for s in sections if s.meta.get("source")=='vision'),    # (think)
		},
	)


# ---- load_pptx_vision  (Pattern A) -----------------------------------------
# Pattern is identical to docx vision: walk slides, then walk picture shapes.
# All blanks resolved as worked examples so you study the parallel — typing
# the same pattern a second time wouldn't add learning. Spend your attention
# on the differences (shape iteration, MSO_SHAPE_TYPE filter, slide_idx).

# _iter_slide_sections_for_slide: single-slide version of _iter_slide_sections.
# Provided in full so vision can process slides one at a time without rewriting native logic.
def _iter_slide_sections_for_slide(slide, idx: int) -> Iterator[ParsedSection]:
	title_shape = slide.shapes.title
	title = title_shape.text_frame.text.strip() if title_shape else ""
	body_parts = [
		shape.text_frame.text.strip()
		for shape in slide.shapes
		if shape != title_shape and shape.has_text_frame and shape.text_frame.text.strip()
	]
	body = "\n".join(body_parts)
	notes = None
	if slide.has_notes_slide:
		notes = slide.notes_slide.notes_text_frame.text.strip() or None
	if not (title or body or notes):
		return
	yield ParsedSection(
		text=body or title, section="slide", slide_idx=idx,
		heading_path=[title] if title else [], speaker_notes=notes,
	)


# ---- WORKED EXAMPLE (no blanks — study the parallel to docx vision) --------
def _iter_pptx_vision_sections(prs, vision_client) -> Iterator[ParsedSection]:
	"""Native text per slide, then one image-section per picture shape (Pattern A)."""
	from pptx.enum.shapes import MSO_SHAPE_TYPE

	for idx, slide in enumerate(prs.slides, start=1):
		yield from _iter_slide_sections_for_slide(slide, idx)

		if vision_client is None:
			continue
		for shape in slide.shapes:
			if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
				continue
			# GIVEN: pptx picture shapes expose raw bytes at shape.image.blob.
			# Same .blob convention as docx; different attribute path.
			image_bytes = shape.image.blob
			description = vision_client(image_bytes, _VISION_PROMPT)
			if description.strip():
				yield ParsedSection(
					text=description,
					section="figure",                # given — same answer as docx 1C
					slide_idx=idx,
					meta={"source": "vision"},
				)


def load_pptx_vision(path: Path, vision_client=None) -> ParsedDocument:
	"""Approach vision (1C): native pptx text + vision descriptions of picture shapes."""
	prs = Presentation(path)
	sections = list(_iter_pptx_vision_sections(prs, vision_client))
	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="pptx",
		title=path.stem,
		sections=sections,
		meta={
			"image_count":    _count_pptx_images(prs),
			# GIVEN — identical pattern to load_docx_1c above. Don't type it twice.
			"vision_applied": sum(1 for s in sections if s.meta.get("source") == "vision"),
		},
	)


# ---- load_pdf_vision  (Pattern B + figure overlay) -------------------------
# Two independent triggers can fire per page:
#   Trigger 1 (replace): native too thin   → vision REPLACES native.
#   Trigger 2 (overlay): native fine BUT page is figure-shaped → vision ADDS a
#                        second "figure" section without dropping native text.
#
# meta["source"] tags each yielded section so downstream eval can compare modes:
#   "native"          | "vision" (replacement) | "vision_overlay" (augmentation)


# ---- GIVEN helpers (study, then use) ---------------------------------------
# DRY reuse of the /Image-XObject walker from load_pdf_1b — but as `any()`
# instead of `sum()` so we short-circuit (one match is enough; no need to count).
def _page_has_image_xobject(page) -> bool:
	"""True if this PDF page references any /Image XObject."""
	resources = page.get("/Resources")
	if not resources or "/XObject" not in resources:
		return False
	xobj = resources["/XObject"].get_object()
	return any(obj.get_object().get("/Subtype") == "/Image" for obj in xobj.values())


# Extracted from the previous inline pdf2image + vision call (DRY: now called
# from two branches). first_page=last_page=idx renders exactly one page.
def _vision_describe_page(path: Path, idx: int, vision_client) -> str:
	"""Render one PDF page (1-indexed) to JPEG and ask vision to describe it."""
	from pdf2image import convert_from_path                       # pip install pdf2image + poppler
	images = convert_from_path(str(path), first_page=idx, last_page=idx)
	return vision_client(_pil_to_bytes(images[0]), _VISION_PROMPT).strip()


# ---- Scaffold 1 of 2 -------------------------------------------------------
def _has_figure_signal(page, text: str) -> bool:
	"""True if the page is figure-shaped: has BOTH an embedded image AND text
	mentioning 'Figure N' / 'Fig. N' / 'Table N'.

	Why both? Either alone is too noisy:
	  - image-only:    logos, watermarks, borders trigger false positives.
	  - caption-only:  "List of Figures" pages, body references to figures
	                   that physically live on a different page.
	The conjunction filters most of those out.
	"""
	has_image   = _page_has_image_xobject(page)
	has_caption = bool(_FIGURE_CAPTION_RE.search(text))

	# Compose the two predicates. We need BOTH to be true to call vision.
	# Pick the boolean operator that short-circuits — if has_image is False, we
	# skip the regex check entirely (cheaper on text-heavy pages with no images).
	return has_image and has_caption                         # (easy) algorithmic — boolean operator


# ---- Scaffold 2 of 2 -------------------------------------------------------
def _iter_pdf_vision_sections(path: Path, reader, vision_client) -> Iterator[ParsedSection]:
	"""Hybrid Pattern B + figure overlay.

	For each page:
	  Branch A (native primary): no vision OR text is sufficient.
	    - Yield the native page section.
	    - Then: if vision is available AND the page is figure-shaped,
	      ALSO yield a vision_overlay "figure" section (augmentation).
	  Branch B (vision replacement): native too thin AND vision available.
	    - Yield a single vision-derived "page" section.
	"""
	for idx, page in enumerate(reader.pages, start=1):
		text = (page.extract_text() or "").strip()

		text_sufficient = len(text.split()) >= _VISION_THRESHOLD
		can_use_vision  = vision_client is not None

		# ---- Branch A: native is primary -----------------------------------
		if not can_use_vision or text_sufficient:
			if text:
				yield ParsedSection(text=text, section="page", page_idx=idx,
				                    meta={"source": "native"})

			# Overlay trigger: compose the conditions for "we want vision in
			# addition to native on this page."
			# Two facts must hold:
			#   (1) we actually HAVE a vision client to call
			#   (2) the page looks figure-shaped (use the helper above)
			# Combine them with the same short-circuit operator you used in
			# _has_figure_signal — and put the cheaper check first so we skip
			# the helper call when vision isn't available anyway.
			if can_use_vision and _has_figure_signal(page, text):                   # (think) algorithmic — compose 2 conditions
				desc = _vision_describe_page(path, idx, vision_client)
				if desc:
					yield ParsedSection(
						text=desc,
						section="figure",
						page_idx=idx,
						# Hint: name the source mode as a single snake_case label.
						# Background: existing values are "native" (no vision) and
						#             "vision" (vision replaced native). We need a
						#             third value meaning "vision called, but native
						#             was kept too." A two-word compound fits.
						# Answer: "vision_overlay"
						meta={"source": "vision_overlay"},           # (easy) trivia — see comment
					)

		# ---- Branch B: native too thin → vision replaces -------------------
		else:
			desc = _vision_describe_page(path, idx, vision_client)
			if not desc:
				continue
			yield ParsedSection(text=desc, section="page", page_idx=idx,
			                    meta={"source": "vision"})


# ---- Scaffold 2 of 2 (no blanks — vision_applied pattern already learned) -
def load_pdf_vision(path: Path, vision_client=None) -> ParsedDocument:
	"""Approach vision (1C): hybrid native + vision PDF loader."""
	reader = PdfReader(str(path))
	sections = list(_iter_pdf_vision_sections(path, reader, vision_client))
	return ParsedDocument(
		source=str(path.resolve()),
		doc_type="pdf",
		title=path.stem,
		sections=sections,
		meta={
			"image_count":    _count_pdf_images(reader),
			"page_count":     len(reader.pages),
			# GIVEN — same pattern as the docx and pptx 1C loaders.
			"vision_applied": sum(1 for s in sections if s.meta.get("source") == "vision"),
		},
	)


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

LOADERS_BY_APPROACH = {
	"1a": {".docx": load_docx_unstructured, ".pptx": load_pptx_unstructured, ".pdf": load_pdf_unstructured},
	"1b": {".docx": load_docx_native, ".pptx": load_pptx_native, ".pdf": load_pdf_native},
	"1c": {".docx": load_docx_vision, ".pptx": load_pptx_vision, ".pdf": load_pdf_vision},
}

load_docx = load_docx_native   # backward-compatible aliases
load_pptx = load_pptx_native
load_pdf  = load_pdf_native

# Numeric-suffix aliases — keep old names working during transition
load_docx_1a = load_docx_unstructured
load_pptx_1a = load_pptx_unstructured
load_pdf_1a  = load_pdf_unstructured
load_docx_1b = load_docx_native
load_pptx_1b = load_pptx_native
load_pdf_1b  = load_pdf_native
load_docx_1c = load_docx_vision
load_pptx_1c = load_pptx_vision
load_pdf_1c  = load_pdf_vision


def load_one(path: Path, approach: str = "1b") -> ParsedDocument:
	"""Load one file with the chosen approach."""
	normalized = _normalize_approach(approach)
	loaders = LOADERS_BY_APPROACH[normalized]
	ext = path.suffix.lower()
	if ext not in loaders:
		raise ValueError(f"Unsupported file type: {ext}")
	return loaders[ext](path)


def load_many(paths: list[Path], approach: str = "1b") -> Iterator[ParsedDocument]:
	"""Yield ParsedDocuments for a list of paths."""
	for p in paths:
		yield load_one(p, approach=approach)


def write_jsonl(docs: list[ParsedDocument], out_path: Path) -> None:
	"""Write one JSON object per line to out_path."""
	with out_path.open("w", encoding="utf-8") as fh:
		for doc in docs:
			fh.write(json.dumps(doc.to_dict(), ensure_ascii=False) + "\n")


# ---------------------------------------------------------------------------
# Smoke test  (run: python -m src.ingestion.loaders)
# ---------------------------------------------------------------------------

if __name__ == "__main__":
	from pathlib import Path
	sample = Path("tests/fixtures/sample.docx")   # swap in any .docx you have
	doc = load_docx_native(sample)
	print(f"doc_type : {doc.doc_type}")
	print(f"sections : {doc.section_count}")
	print(f"meta     : {doc.meta}")
	print(f"preview  : {doc.full_text[:120]}")
	# expected shape (numbers will vary by file):
	# doc_type : docx
	# sections : 14
	# meta     : {'image_count': 2}
	# preview  : Introduction\n\nThis thesis investigates...


# ---- HINTS (uncover only if stuck > 5 min) ----
# Scaffold 1 — _heading_level
#   Step 1: str.split() with no args splits on any whitespace.
#   Step 2: the str method is .isdigit().
#   Step 3: the built-in is int().
# Scaffold 2 — _iter_doc_sections
#   is_heading: docx style names look like "heading 1", "heading 2" — what word is always there?
#   slice math: if a level-2 heading appears, you want to keep ancestors at depth < 2,
#               i.e. the first (level - 1) items of the stack. So  heading_stack[:level - 1].
#   non-heading section name: matches what the original code used.
# Scaffold 3 — load_docx_native
#   No blanks. If unsure, re-read Step 2: we call crawl.stat_docx instead of re-implementing
#   the image-counting loop. DRY = "Don't Repeat Yourself".


# ---- Reflection questions ----
# Q1 (_heading_level): Why guard with `parts and parts[-1].isdigit()` rather than just
#    `parts[-1].isdigit()`? Name one input that crashes the simpler version.
#    (One-word answer involves: empty.)
#
# Q2 (_iter_doc_sections): We do  heading_stack = heading_stack[:level - 1]  instead of
#    popping items in a while-loop. Both work. Name one advantage of the slice approach
#    AND one situation where the while-loop version would actually be preferable.
#
# Q3 (load_docx_native): By calling crawl.stat_docx(path) we open the same .docx file TWICE
#    (once here, once inside stat_docx). What's the cost of that as the corpus grows,
#    and how would you fix it without losing the "single source of truth" benefit?


# ---- Reflection questions for load_pdf_native ----
# Q4: We skip pages where extract_text() returns empty. For a scanned PhD thesis, EVERY
#     page would be empty. What metadata would you add to ParsedDocument.meta so a
#     downstream OCR pass knows "this PDF needs OCR" without re-opening the file?
#
# Q5: extract_text() can sometimes return text in a weird order (multi-column layouts,
#     footnotes interleaved with body, etc.). Name one concrete failure mode of the
#     "one-section-per-page" choice for a multi-column research PDF, and how you might
#     detect it during ingestion (not at retrieval time).


# ---- HINTS for unstructured loaders (uncover only if stuck > 5 min) ----
# _unstructured_section_label:
#   The heading category string in unstructured is exactly "Title" (capital T).
#
# _iter_docx_unstructured_sections:
#   Category to skip: "Image"
#   last_title update: last_title = text
#   heading heading_path: [text]          (heading includes itself — same as native)
#   normal  heading_path: [last_title] if last_title else []
#
# load_docx_unstructured:
#   partition_docx(filename=str(path))
#   image_count: sum(1 for el in elements if el.category == "Image")
#
# _iter_pptx_unstructured_sections / _iter_pdf_unstructured_sections:
#   page/slide index: el.metadata.page_number   (may be None — fine to pass to ParsedSection)
#
# load_pptx_unstructured / load_pdf_unstructured:
#   partition_pptx(filename=str(path)) / partition_pdf(filename=str(path))
#   image_count: same sum(... el.category == "Image") pattern
#
# load_pdf_unstructured — page_count (stretch):
#   len({el.metadata.page_number for el in elements if el.metadata.page_number is not None})


# ---- Reflection questions for unstructured loaders ----
# Q6 (_unstructured_section_label): unstructured has richer categories than just "Title" and
#    "NarrativeText" — e.g. "Table", "FigureCaption", "ListItem". Our label function collapses
#    all non-Title elements to "normal". Name one category where that loses useful structure,
#    and describe the ParsedSection field you'd add to preserve it.
#
# Q7 (load_docx_unstructured vs load_docx_native): The unstructured loader doesn't reconstruct a full heading stack
#    (only the immediately preceding Title, not multi-level ancestry). Under what corpus
#    conditions does that matter for retrieval quality, and how would you fix it?
#
# Q8 (load_pdf_unstructured page_count): We compute page_count from metadata rather than from
#    len(reader.pages) as in native. What could cause the two counts to disagree, and which
#    is more reliable for a scanned PDF where some pages have no extractable elements?


# ---- HINTS for vision loaders ----
# Trivia blanks (JPEG, .blob, first_page/last_page) have hint+background+answer
# inline above each blank — see the function bodies. Below: hints for the THREE
# algorithmic blanks. These ones live at the bottom on purpose — try to reason
# from surrounding code first; uncover only after >5 min stuck.
#
# Algo Blank 1 — _iter_docx_vision_sections, section="________":
#   Concept:    name the unit you're yielding (consistency with docx/pptx/pdf labels)
#   Hint:       the unit here is a vision-described embedded image. Reach for
#               a single noun that describes the visual asset, not the medium.
#   Answer:     "figure"
#
# Algo Blank 2 — load_docx_vision, vision_applied = ____:
#   Concept:    comprehension + filter idiom — count items in a list matching a predicate
#   Hint:       every section that used vision tagged itself with meta["source"] = "vision".
#               Walk `sections`, filter by that meta key, count the matches.
#   Answer:     sum(1 for s in sections if s.meta.get("source") == "vision")
#               (the .get("source") form is safer than s.meta["source"] —
#                native sections may not have the key set.)
#
# Algo Blank 3 — _has_figure_signal, return has_image ____ has_caption:
#   Concept:    boolean operator choice (AND vs OR) for short-circuit composition
#   Hint:       we want vision to fire ONLY when both signals agree. Either alone is
#               too noisy (logos trigger image-only; "List of Figures" triggers caption-only).
#               Pick the operator that requires BOTH operands true.
#   Bonus:      `and` short-circuits — if has_image is False, the regex check is skipped.
#   Answer:     and
#
# Algo Blank 4 — _iter_pdf_vision_sections, overlay condition `if ____:`:
#   Concept:    compose two predicates with the same short-circuit operator
#   Hint:       you want to fire the overlay ONLY when (a) we actually have a vision
#               client AND (b) the page is figure-shaped (use _has_figure_signal).
#               Order the cheaper check first so we skip the helper call when there's
#               no client to call anyway.
#   Answer:     can_use_vision and _has_figure_signal(page, text)


# ---- Reflection questions for vision loaders ----
# Q9 (Pattern A vs B): A searchable literature PDF has plenty of text but its figures
#    (e.g. cycle-life plots, EIS spectra) carry the key findings. Pattern B as written
#    never triggers vision on those pages because word count >= threshold. Sketch one
#    approach that would describe figures even on text-rich pages without re-rendering
#    every page through the vision model.
#
# Q10 (vision_applied meta): We store vision_applied as a count. Name one downstream
#    use case where storing the list of page_idx values that used vision would be more
#    useful than the raw integer.
#
# Q11 (_pil_to_bytes): We hardcode JPEG. Name one class of document where JPEG could
#    silently degrade quality compared to PNG, and explain why (think about what lossy
#    compression does to fine-grained detail).
